"""
Document Indexer Module - Document processing and indexing pipeline

This module handles the complete document processing workflow:
- File discovery from configured directories
- Text extraction from multiple formats (PDF, DOCX, PPTX, TXT, MD)
- Character-based text chunking with overlap
- Batch processing with progress tracking
- Integration with vector store for embedding generation
"""

import hashlib
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Any
from dataclasses import dataclass
from datetime import datetime

from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TaskProgressColumn

from src.config import Config
from src.vector_store import VectorStore, Document, generate_document_id


# ============================================================================
# Constants
# ============================================================================

console = Console()


# ============================================================================
# Data Models
# ============================================================================

@dataclass
class Chunk:
    """
    Text chunk with metadata.
    
    Attributes:
        text: The chunk text content
        index: 0-based position in document
        page: Optional page number (for PDFs)
    """
    text: str
    index: int
    page: Optional[int] = None


@dataclass
class ExtractedText:
    """
    Extracted text from a document with page information.
    
    Attributes:
        text: Complete extracted text
        pages: Optional list of (page_number, page_text) tuples for PDFs
    """
    text: str
    pages: Optional[List[Tuple[int, str]]] = None


@dataclass
class IndexStats:
    """
    Statistics from an indexing operation.
    
    Attributes:
        files_processed: Number of files successfully processed
        files_failed: Number of files that failed processing
        files_skipped: Number of files skipped (no extractable text)
        chunks_created: Total number of chunks created
        duration_seconds: Time taken for indexing
        failed_files: List of file paths that failed
    """
    files_processed: int = 0
    files_failed: int = 0
    files_skipped: int = 0
    chunks_created: int = 0
    duration_seconds: float = 0.0
    failed_files: List[Path] = None
    
    def __post_init__(self):
        if self.failed_files is None:
            self.failed_files = []


# ============================================================================
# File Discovery
# ============================================================================

def discover_files(
    directories: List[str],
    file_types: List[str]
) -> List[Path]:
    """
    Discover all files matching the specified types in directories.
    
    Recursively scans directories and filters by file extension.
    Skips hidden directories and system folders.
    
    Args:
        directories: List of directory paths to scan
        file_types: List of file extensions to include (e.g., [".pdf", ".txt"])
        
    Returns:
        List of Path objects for matching files
    """
    discovered_files = []
    valid_extensions = set(ext.lower() for ext in file_types)
    
    # Directories to skip
    skip_dirs = {".git", ".venv", "venv", "node_modules", "__pycache__", 
                 ".pytest_cache", "htmlcov", "data"}
    
    for directory_str in directories:
        directory = Path(directory_str).expanduser().resolve()
        
        if not directory.exists() or not directory.is_dir():
            console.print(f"[yellow]Warning: Directory not found: {directory}[/yellow]")
            continue
        
        try:
            # Recursively find files
            for item in directory.rglob("*"):
                # Skip if in excluded directory
                if any(excluded in item.parts for excluded in skip_dirs):
                    continue
                
                # Skip hidden files and directories
                if any(part.startswith(".") for part in item.parts[len(directory.parts):]):
                    continue
                
                # Check if file matches extension
                if item.is_file() and item.suffix.lower() in valid_extensions:
                    discovered_files.append(item)
                    
        except PermissionError:
            console.print(f"[yellow]Warning: Permission denied accessing {directory}[/yellow]")
            continue
    
    return discovered_files


# ============================================================================
# Text Extraction
# ============================================================================

def extract_text_from_pdf(file_path: Path) -> ExtractedText:
    """
    Extract text from PDF file using PyPDF2.
    
    Args:
        file_path: Path to PDF file
        
    Returns:
        ExtractedText object with text and page information
        
    Raises:
        Exception: If PDF cannot be read
    """
    try:
        import PyPDF2
        
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        pages.append((page_num, page_text))
                except Exception as e:
                    console.print(f"[yellow]Warning: Failed to extract page {page_num} from {file_path.name}: {e}[/yellow]")
                    continue
            
            # Combine all pages
            full_text = "\n\n".join(text for _, text in pages)
            
            return ExtractedText(text=full_text, pages=pages)
            
    except Exception as e:
        raise Exception(f"Failed to extract PDF: {e}")


def extract_text_from_docx(file_path: Path) -> ExtractedText:
    """
    Extract text from DOCX file using python-docx.
    
    Args:
        file_path: Path to DOCX file
        
    Returns:
        ExtractedText object with text
        
    Raises:
        Exception: If DOCX cannot be read
    """
    try:
        from docx import Document as DocxDocument
        
        doc = DocxDocument(file_path)
        paragraphs = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        paragraphs.append(text)
        
        full_text = "\n\n".join(paragraphs)
        
        return ExtractedText(text=full_text)
        
    except Exception as e:
        raise Exception(f"Failed to extract DOCX: {e}")


def extract_text_from_pptx(file_path: Path) -> ExtractedText:
    """
    Extract text from PPTX file using python-pptx.
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        ExtractedText object with text
        
    Raises:
        Exception: If PPTX cannot be read
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        slides_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if slide_content:
                slides_text.append(f"Slide {slide_num}:\n" + "\n".join(slide_content))
        
        full_text = "\n\n".join(slides_text)
        
        return ExtractedText(text=full_text)
        
    except Exception as e:
        raise Exception(f"Failed to extract PPTX: {e}")


def extract_text_from_plain(file_path: Path) -> ExtractedText:
    """
    Extract text from plain text files (TXT, MD).
    
    Args:
        file_path: Path to text file
        
    Returns:
        ExtractedText object with text
        
    Raises:
        Exception: If file cannot be read
    """
    try:
        # Try UTF-8 first, fall back to other encodings
        encodings = ["utf-8", "latin-1", "cp1252"]
        
        for encoding in encodings:
            try:
                text = file_path.read_text(encoding=encoding)
                return ExtractedText(text=text)
            except UnicodeDecodeError:
                continue
        
        # If all fail, read with error handling
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        return ExtractedText(text=text)
        
    except Exception as e:
        raise Exception(f"Failed to read text file: {e}")


def extract_text(file_path: Path) -> ExtractedText:
    """
    Extract text from file based on extension.
    
    Routes to appropriate extraction function based on file type.
    
    Args:
        file_path: Path to file
        
    Returns:
        ExtractedText object
        
    Raises:
        ValueError: If file type not supported
        Exception: If extraction fails
    """
    suffix = file_path.suffix.lower()
    
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    elif suffix == ".docx":
        return extract_text_from_docx(file_path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    elif suffix in [".txt", ".md"]:
        return extract_text_from_plain(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


# ============================================================================
# Text Chunking
# ============================================================================

def chunk_text(
    text: str,
    chunk_size: int,
    chunk_overlap: int,
    pages: Optional[List[Tuple[int, str]]] = None
) -> List[Chunk]:
    """
    Split text into overlapping chunks.
    
    Uses simple character-based chunking. Does not attempt to preserve
    sentence boundaries (MVP approach).
    
    Args:
        text: Text to chunk
        chunk_size: Maximum characters per chunk
        chunk_overlap: Number of overlapping characters between chunks
        pages: Optional page information for tracking page numbers
        
    Returns:
        List of Chunk objects
    """
    if not text.strip():
        return []
    
    chunks = []
    start = 0
    index = 0
    
    while start < len(text):
        # Extract chunk
        end = start + chunk_size
        chunk_text = text[start:end]
        
        # Skip empty chunks
        if not chunk_text.strip():
            start = end
            continue
        
        # Determine page number if pages provided
        page_num = None
        if pages:
            # Find which page this chunk belongs to
            char_count = 0
            for page_n, page_text in pages:
                char_count += len(page_text) + 2  # +2 for \n\n separator
                if start < char_count:
                    page_num = page_n
                    break
        
        chunks.append(Chunk(
            text=chunk_text.strip(),
            index=index,
            page=page_num
        ))
        
        index += 1
        
        # Move to next chunk with overlap
        start = end - chunk_overlap
        
        # Avoid infinite loop if overlap >= chunk_size
        if start <= end - chunk_size:
            start = end
    
    return chunks


# ============================================================================
# Document Processing
# ============================================================================

def process_document(
    file_path: Path,
    config: Config
) -> List[Document]:
    """
    Process a single document: extract text, chunk, and prepare for indexing.
    
    Args:
        file_path: Path to document
        config: Application configuration
        
    Returns:
        List of Document objects ready for vector store
        
    Raises:
        Exception: If processing fails
    """
    # Extract text
    extracted = extract_text(file_path)
    
    # Check if we got any text
    if not extracted.text.strip():
        raise Exception("No text content extracted")
    
    # Chunk text
    chunks = chunk_text(
        text=extracted.text,
        chunk_size=config.chunk_size,
        chunk_overlap=config.chunk_overlap,
        pages=extracted.pages
    )
    
    if not chunks:
        raise Exception("No chunks created")
    
    # Get file metadata
    file_stats = file_path.stat()
    
    # Create Document objects
    documents = []
    for chunk in chunks:
        doc_id = generate_document_id(file_path, chunk.index)
        
        metadata = {
            "source": str(file_path.resolve()),
            "filename": file_path.name,
            "chunk_index": chunk.index,
            "total_chunks": len(chunks),
            "file_type": file_path.suffix.lower(),
            "indexed_at": datetime.now().isoformat(),
            "file_size": file_stats.st_size,
            "file_modified": datetime.fromtimestamp(file_stats.st_mtime).isoformat()
        }
        
        # Add page number if available
        if chunk.page is not None:
            metadata["page"] = chunk.page
        
        documents.append(Document(
            id=doc_id,
            text=chunk.text,
            metadata=metadata
        ))
    
    return documents


# ============================================================================
# Document Indexer
# ============================================================================

class DocumentIndexer:
    """
    Orchestrates the document indexing pipeline.
    
    Handles file discovery, processing, and storage in vector store
    with progress tracking and error handling.
    """
    
    def __init__(self, config: Config, vector_store: VectorStore):
        """
        Initialize indexer.
        
        Args:
            config: Application configuration
            vector_store: Initialized vector store instance
        """
        self.config = config
        self.vector_store = vector_store
    
    def index_all(self, show_progress: bool = True) -> IndexStats:
        """
        Index all documents from configured directories.
        
        Main orchestration method that:
        1. Discovers files
        2. Processes each file (extract → chunk)
        3. Batches documents for vector store
        4. Tracks progress and errors
        
        Args:
            show_progress: Whether to display progress bar
            
        Returns:
            IndexStats with results
        """
        stats = IndexStats()
        start_time = datetime.now()
        
        # Discover files
        console.print("[cyan]Discovering documents...[/cyan]")
        files = discover_files(self.config.directories, self.config.file_types)
        
        if not files:
            console.print("[yellow]No documents found to index[/yellow]")
            return stats
        
        console.print(f"[green]Found {len(files)} files to process[/green]\n")
        
        # Process files with progress bar
        if show_progress:
            with Progress(
                SpinnerColumn(),
                TextColumn("[progress.description]{task.description}"),
                BarColumn(),
                TaskProgressColumn(),
                console=console
            ) as progress:
                task = progress.add_task("[cyan]Indexing documents...", total=len(files))
                
                for file_path in files:
                    self._process_file_with_progress(file_path, stats, progress, task)
                    progress.advance(task)
        else:
            # Process without progress bar
            for file_path in files:
                self._process_file(file_path, stats)
        
        # Calculate duration
        end_time = datetime.now()
        stats.duration_seconds = (end_time - start_time).total_seconds()
        
        # Display summary
        console.print()
        console.print("[bold green]Indexing complete![/bold green]")
        console.print(f"  Files processed: {stats.files_processed}")
        console.print(f"  Chunks created: {stats.chunks_created}")
        console.print(f"  Files failed: {stats.files_failed}")
        console.print(f"  Files skipped: {stats.files_skipped}")
        console.print(f"  Duration: {stats.duration_seconds:.2f}s")
        
        if stats.failed_files:
            console.print("\n[yellow]Failed files:[/yellow]")
            for failed in stats.failed_files[:10]:  # Show first 10
                console.print(f"  - {failed}")
            if len(stats.failed_files) > 10:
                console.print(f"  ... and {len(stats.failed_files) - 10} more")
        
        return stats
    
    def _process_file(self, file_path: Path, stats: IndexStats) -> None:
        """
        Process a single file without progress tracking.
        
        Args:
            file_path: Path to file
            stats: Statistics object to update
        """
        try:
            documents = process_document(file_path, self.config)
            
            if documents:
                # Add to vector store
                self.vector_store.add_documents(documents)
                
                stats.files_processed += 1
                stats.chunks_created += len(documents)
            else:
                stats.files_skipped += 1
                
        except Exception as e:
            console.print(f"[yellow]Warning: Failed to process {file_path.name}: {e}[/yellow]")
            stats.files_failed += 1
            stats.failed_files.append(file_path)
    
    def _process_file_with_progress(
        self,
        file_path: Path,
        stats: IndexStats,
        progress: Progress,
        task
    ) -> None:
        """
        Process a single file with progress bar update.
        
        Args:
            file_path: Path to file
            stats: Statistics object to update
            progress: Rich Progress instance
            task: Progress task ID
        """
        # Update progress description
        progress.update(task, description=f"[cyan]Processing {file_path.name[:40]}...")
        
        # Process file
        self._process_file(file_path, stats)


# ============================================================================
# Module Testing
# ============================================================================

if __name__ == "__main__":
    # For testing indexer independently
    from src.config import get_default_config
    
    config = get_default_config()
    config.directories = ["."]
    
    # Test file discovery
    files = discover_files(config.directories, config.file_types)
    print(f"Found {len(files)} files")
    
    # Test text extraction on a sample file if available
    if files:
        sample_file = files[0]
        print(f"\nTesting extraction on: {sample_file}")
        try:
            extracted = extract_text(sample_file)
            print(f"Extracted {len(extracted.text)} characters")
            
            # Test chunking
            chunks = chunk_text(extracted.text, 500, 100)
            print(f"Created {len(chunks)} chunks")
        except Exception as e:
            print(f"Error: {e}")
