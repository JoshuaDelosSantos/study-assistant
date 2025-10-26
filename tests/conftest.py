"""
Pytest Configuration and Shared Fixtures

This module contains pytest fixtures that are shared across all test files.
Fixtures defined here are automatically available in all test modules.
"""

import os
import tempfile
from pathlib import Path
from typing import Generator
import pytest

from src.config import Config


# ============================================================================
# Temporary Directory Fixtures
# ============================================================================

@pytest.fixture
def temp_dir() -> Generator[Path, None, None]:
    """
    Create a temporary directory for test files.
    
    Yields:
        Path to temporary directory
        
    Cleanup:
        Automatically removed after test
    """
    with tempfile.TemporaryDirectory() as tmp_dir:
        yield Path(tmp_dir)


@pytest.fixture
def temp_config_dir(temp_dir: Path, monkeypatch) -> Path:
    """
    Create a temporary directory and set it as the config location.
    
    Args:
        temp_dir: Temporary directory fixture
        monkeypatch: Pytest monkeypatch fixture
        
    Returns:
        Path to temporary config directory
    """
    # Monkey patch the config module to use temp directory
    import src.config as config_module
    
    config_file = temp_dir / "config.yaml"
    data_dir = temp_dir / "data"
    data_dir.mkdir()
    
    monkeypatch.setattr(config_module, "CONFIG_FILE", config_file)
    monkeypatch.setattr(config_module, "DATA_DIR", data_dir)
    
    return temp_dir


# ============================================================================
# Sample Directory Structure Fixtures
# ============================================================================

@pytest.fixture
def sample_study_structure(temp_dir: Path) -> dict[str, Path]:
    """
    Create a sample study directory structure for testing.
    
    Creates:
        temp_dir/
        ├── CS101/
        │   ├── lecture1.pdf (empty placeholder)
        │   ├── notes.txt
        │   └── assignments/
        │       └── hw1.docx (empty placeholder)
        ├── Math202/
        │   ├── textbook.pdf (empty placeholder)
        │   └── exercises.txt
        └── EmptyFolder/
    
    Returns:
        Dictionary mapping directory names to paths
    """
    # Create directories
    cs101 = temp_dir / "CS101"
    cs101.mkdir()
    
    math202 = temp_dir / "Math202"
    math202.mkdir()
    
    empty = temp_dir / "EmptyFolder"
    empty.mkdir()
    
    assignments = cs101 / "assignments"
    assignments.mkdir()
    
    # Create files
    (cs101 / "lecture1.pdf").write_text("PDF content placeholder")
    (cs101 / "notes.txt").write_text("CS101 lecture notes\nPython basics")
    (assignments / "hw1.docx").write_text("DOCX content placeholder")
    
    (math202 / "textbook.pdf").write_text("Math textbook placeholder")
    (math202 / "exercises.txt").write_text("Math exercises\nCalculus problems")
    
    return {
        "root": temp_dir,
        "cs101": cs101,
        "math202": math202,
        "empty": empty,
        "assignments": assignments,
    }


# ============================================================================
# Configuration Fixtures
# ============================================================================

@pytest.fixture
def default_config() -> Config:
    """
    Create a default configuration for testing.
    
    Returns:
        Config instance with default values
    """
    return Config()


@pytest.fixture
def valid_config(temp_dir: Path) -> Config:
    """
    Create a valid configuration for testing.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Valid Config instance
    """
    config = Config(
        llm_provider="openai",
        llm_model="gpt-4o",
        llm_api_key="sk-test-key-12345",
        directories=[str(temp_dir)],
    )
    return config


@pytest.fixture
def config_dict() -> dict:
    """
    Create a configuration dictionary for testing.
    
    Returns:
        Dictionary with valid configuration
    """
    return {
        "llm_provider": "openai",
        "llm_model": "gpt-4o",
        "llm_api_key": "sk-test-key-12345",
        "llm_temperature": 0.7,
        "llm_max_tokens": 2000,
        "embedding_provider": "sentence-transformers",
        "embedding_model": "all-MiniLM-L6-v2",
        "directories": ["/test/dir1", "/test/dir2"],
        "chunk_size": 1000,
        "chunk_overlap": 200,
        "file_types": [".pdf", ".docx", ".txt"],
        "top_k": 5,
        "similarity_threshold": 0.7,
        "track_tokens": True,
        "show_sources": True,
        "verbose": False,
    }


# ============================================================================
# Environment Fixtures
# ============================================================================

@pytest.fixture
def clean_env(monkeypatch) -> None:
    """
    Clean environment variables that might interfere with tests.
    
    Args:
        monkeypatch: Pytest monkeypatch fixture
    """
    # Remove API keys from environment
    env_vars = ["OPENAI_API_KEY", "GOOGLE_API_KEY", "GEMINI_API_KEY"]
    for var in env_vars:
        monkeypatch.delenv(var, raising=False)


@pytest.fixture
def mock_openai_env(monkeypatch) -> str:
    """
    Set up mock OpenAI API key in environment.
    
    Args:
        monkeypatch: Pytest monkeypatch fixture
        
    Returns:
        The mock API key
    """
    api_key = "sk-test-openai-key"
    monkeypatch.setenv("OPENAI_API_KEY", api_key)
    return api_key


@pytest.fixture
def mock_gemini_env(monkeypatch) -> str:
    """
    Set up mock Gemini API key in environment.
    
    Args:
        monkeypatch: Pytest monkeypatch fixture
        
    Returns:
        The mock API key
    """
    api_key = "test-gemini-key"
    monkeypatch.setenv("GOOGLE_API_KEY", api_key)
    return api_key


# ============================================================================
# Mock Fixtures
# ============================================================================

@pytest.fixture
def mock_console(mocker):
    """
    Mock Rich console for testing CLI output.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock console object
    """
    return mocker.patch("src.cli.console")


@pytest.fixture
def mock_prompt(mocker):
    """
    Mock Rich Prompt for testing user input.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock Prompt object
    """
    return mocker.patch("src.setup_wizard.Prompt")


@pytest.fixture
def mock_confirm(mocker):
    """
    Mock Rich Confirm for testing yes/no prompts.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock Confirm object
    """
    return mocker.patch("src.setup_wizard.Confirm")


# ============================================================================
# Vector Store Fixtures (Phase 2)
# ============================================================================

@pytest.fixture
def temp_data_dir(temp_dir: Path, monkeypatch) -> Path:
    """
    Create a temporary directory and set it as the data location.
    
    Args:
        temp_dir: Temporary directory fixture
        monkeypatch: Pytest monkeypatch fixture
        
    Returns:
        Path to temporary data directory
    """
    import src.config as config_module
    import src.vector_store as vector_store_module
    
    data_dir = temp_dir / "data"
    data_dir.mkdir()
    
    monkeypatch.setattr(config_module, "DATA_DIR", data_dir)
    monkeypatch.setattr(vector_store_module, "DATA_DIR", data_dir)
    
    return data_dir


@pytest.fixture
def mock_embedding_function(mocker):
    """
    Mock ChromaDB embedding function for fast tests.
    
    Returns a mock that generates random 384-dimensional vectors.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock embedding function
    """
    def mock_embed(texts):
        """Generate fake embeddings."""
        import random
        return [[random.random() for _ in range(384)] for _ in texts]
    
    mock_fn = mocker.MagicMock()
    mock_fn.__call__ = mock_embed
    
    return mock_fn


@pytest.fixture
def sample_pdf_file(temp_dir: Path) -> Path:
    """
    Create a simple PDF file for testing.
    
    Uses reportlab if available, otherwise creates a placeholder.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created PDF file
    """
    pdf_path = temp_dir / "test_document.pdf"
    
    try:
        # Try to create a real PDF with reportlab
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        c.drawString(100, 750, "Test PDF Document")
        c.drawString(100, 730, "This is a test document for the study assistant.")
        c.drawString(100, 710, "It contains multiple lines of text.")
        c.showPage()
        c.drawString(100, 750, "Page 2")
        c.drawString(100, 730, "This is the second page of the test document.")
        c.showPage()
        c.save()
        
    except ImportError:
        # Reportlab not available, create minimal placeholder
        # This won't be readable by PyPDF2, but can be used for file discovery tests
        pdf_path.write_bytes(b"%PDF-1.4\n%\xE2\xE3\xCF\xD3\n")
    
    return pdf_path


@pytest.fixture
def sample_docx_file(temp_dir: Path) -> Path:
    """
    Create a simple DOCX file for testing.
    
    Uses python-docx if available, otherwise creates a placeholder.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created DOCX file
    """
    docx_path = temp_dir / "test_document.docx"
    
    try:
        from docx import Document
        
        doc = Document()
        doc.add_heading("Test DOCX Document", 0)
        doc.add_paragraph("This is a test document for the study assistant.")
        doc.add_paragraph("It contains multiple paragraphs of text.")
        doc.add_heading("Section 2", level=1)
        doc.add_paragraph("This is the second section with more content.")
        doc.save(str(docx_path))
        
    except ImportError:
        # python-docx not available, create placeholder
        docx_path.write_text("DOCX placeholder content")
    
    return docx_path


@pytest.fixture
def sample_pptx_file(temp_dir: Path) -> Path:
    """
    Create a simple PPTX file for testing.
    
    Uses python-pptx if available, otherwise creates a placeholder.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created PPTX file
    """
    pptx_path = temp_dir / "test_presentation.pptx"
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "A sample presentation for testing"
        
        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide2.shapes.title
        content = slide2.placeholders[1]
        title.text = "Slide 2"
        content.text = "This is the content of the second slide."
        
        prs.save(str(pptx_path))
        
    except ImportError:
        # python-pptx not available, create placeholder
        pptx_path.write_text("PPTX placeholder content")
    
    return pptx_path


@pytest.fixture
def sample_txt_file(temp_dir: Path) -> Path:
    """
    Create a simple text file for testing.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created text file
    """
    txt_path = temp_dir / "test_notes.txt"
    txt_path.write_text(
        "Test Text Document\n"
        "\n"
        "This is a test text file for the study assistant.\n"
        "It contains multiple lines of plain text.\n"
        "\n"
        "Section 2\n"
        "This is another section with more content.\n"
        "Plain text files are the simplest to process.\n"
    )
    return txt_path


@pytest.fixture
def sample_md_file(temp_dir: Path) -> Path:
    """
    Create a simple Markdown file for testing.
    
    Args:
        temp_dir: Temporary directory fixture
        
    Returns:
        Path to created Markdown file
    """
    md_path = temp_dir / "test_notes.md"
    md_path.write_text(
        "# Test Markdown Document\n"
        "\n"
        "This is a test markdown file for the study assistant.\n"
        "\n"
        "## Section 1\n"
        "\n"
        "It contains **formatted** text with _emphasis_.\n"
        "\n"
        "## Section 2\n"
        "\n"
        "- Bullet point 1\n"
        "- Bullet point 2\n"
        "- Bullet point 3\n"
    )
    return md_path


@pytest.fixture
def sample_documents_dir(
    temp_dir: Path,
    sample_pdf_file: Path,
    sample_docx_file: Path,
    sample_pptx_file: Path,
    sample_txt_file: Path,
    sample_md_file: Path
) -> Path:
    """
    Create a directory with all sample document types.
    
    Args:
        temp_dir: Temporary directory fixture
        sample_pdf_file: PDF fixture
        sample_docx_file: DOCX fixture
        sample_pptx_file: PPTX fixture
        sample_txt_file: TXT fixture
        sample_md_file: MD fixture
        
    Returns:
        Path to directory containing all samples
    """
    # All files are already in temp_dir, just return it
    return temp_dir


# ============================================================================
# LLM Provider Fixtures (Phase 3)
# ============================================================================

@pytest.fixture
def mock_openai_client(mocker):
    """
    Mock OpenAI client for testing LLM providers.
    
    Returns a mock that simulates OpenAI API responses.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock OpenAI client
    """
    from unittest.mock import Mock
    
    mock_client = Mock()
    mock_response = Mock()
    mock_response.choices = [Mock(message=Mock(content="Mocked OpenAI response"))]
    mock_response.usage = Mock(
        prompt_tokens=10,
        completion_tokens=20,
        total_tokens=30
    )
    mock_response.model = "gpt-4o"
    
    mock_client.chat.completions.create.return_value = mock_response
    
    return mock_client


@pytest.fixture
def mock_gemini_model(mocker):
    """
    Mock Gemini model for testing LLM providers.
    
    Returns a mock that simulates Gemini API responses.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock Gemini model
    """
    from unittest.mock import Mock
    
    mock_model = Mock()
    mock_response = Mock()
    mock_response.text = "Mocked Gemini response"
    mock_response.usage_metadata = Mock(
        prompt_token_count=10,
        candidates_token_count=20,
        total_token_count=30
    )
    
    mock_model.generate_content.return_value = mock_response
    mock_model.count_tokens.return_value = Mock(total_tokens=10)
    
    return mock_model


@pytest.fixture
def mock_tiktoken_encoding(mocker):
    """
    Mock tiktoken encoding for testing token counting.
    
    Returns a mock that simulates tiktoken encoding behavior.
    
    Args:
        mocker: pytest-mock fixture
        
    Returns:
        Mock tiktoken encoding
    """
    from unittest.mock import Mock
    
    mock_encoding = Mock()
    
    def mock_encode(text):
        """Simulate encoding: ~4 chars per token."""
        return [0] * max(1, len(text) // 4)
    
    mock_encoding.encode = mock_encode
    
    return mock_encoding
